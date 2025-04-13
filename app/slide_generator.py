from pptx import Presentation
import os
import requests
from pptx.util import Inches,Pt
import requests
from dotenv import load_dotenv
load_dotenv()
SERP_API_KEY = os.getenv("SERP_API_KEY")

def search_image_serpapi(query):
    params = {
        "q": query,
        "api_key": SERP_API_KEY,
        "engine": "google_images"
    }
    response = requests.get("https://serpapi.com/search", params=params)
    results = response.json()
    if "images_results" in results:
        print(results["images_results"][0]["original"])
        return results["images_results"][0]["original"]
    return None


def download_image(url, save_dir="images"):
    if not url:
        print("No image URL provided.")
        return None

    os.makedirs(save_dir, exist_ok=True)

    image_name = url.split("/")[-1].split("?")[0]
    if not image_name.lower().endswith((".png", ".jpg", ".jpeg")):
        image_name += ".jpg"

    image_path = os.path.join(save_dir, image_name)

    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            with open(image_path, "wb") as f:
                f.write(response.content)
            return image_path
        else:
            print(f"Failed to download image. Status code: {response.status_code}")
    except Exception as e:
        print(f"Image download failed: {e}")

    return None

def create_ppt(slide_json, output_path="lesson_plan.pptx"):
    prs = Presentation()

    for slide in slide_json.get("slides", []):

        # Title Slide
        if slide["type"] == "title":
            ppt_slide = prs.slides.add_slide(prs.slide_layouts[0])
            ppt_slide.shapes.title.text = slide["title"]
            ppt_slide.placeholders[1].top = Inches(2)
            ppt_slide.placeholders[1].height = Inches(2)

            image_path = None
            image_url = search_image_serpapi(slide.get("keyword", ""))
            if image_url:
                image_path = download_image(image_url)
            if image_path:
                title_placeholder = ppt_slide.shapes.title
                title_bottom = title_placeholder.top + title_placeholder.height

                ppt_slide.shapes.add_picture(
                    image_path,
                    Inches(2.5), #horizontal 
                    title_bottom + Inches(0.5), # below the title
                    width=Inches(5.5),
                    height=Inches(3)
                )

        # Content Slide
        elif slide["type"] == "content":
            ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])
            ppt_slide.shapes.title.text = slide["title"]
            content_box = ppt_slide.shapes.placeholders[1]
            content_box.text = ""  
            for point in slide.get("points", []):
                p = content_box.text_frame.add_paragraph()
                run = p.add_run()
                run.text = point
                run.font.size = Pt(18) 


            content_box.left = Inches(0.5)
            content_box.top = Inches(1.5)
            content_box.width = Inches(5.5)
            content_box.height = Inches(4.5)

            image_path = None
            image_url = search_image_serpapi(slide.get("keyword", ""))
            if image_url:
                image_path = download_image(image_url)
            if image_path:
                # Place image to right side of main content
                ppt_slide.shapes.add_picture(
                    image_path, 
                    Inches(6.2), 
                    Inches(3.5), 
                    width=Inches(3), 
                    height=Inches(2.5)
)

        elif slide["type"] == "quiz":
            quiz = slide["quiz"]

            # Quiz Question (without answer)
            question_slide = prs.slides.add_slide(prs.slide_layouts[1])
            question_slide.shapes.title.text = slide["title"]

            if quiz["type"] == "mcq":
                question = f"Question: {quiz['question']}\n\n"
                options = "\n".join([f" {opt}" for opt in quiz["options"]])
                content = question + options

            elif quiz["type"] == "fill-in-the-blank":
                content = f"Fill in the blank: {quiz['question']}"

            question_slide.shapes.placeholders[1].text = content

            # Quiz Answer
            answer_slide = prs.slides.add_slide(prs.slide_layouts[1])
            answer_slide.shapes.title.text = f"{slide['title']} - Answer"

            if quiz["type"] == "mcq":
                answer_content = f"Answer: {quiz['answer']}"
            elif quiz["type"] == "fill-in-the-blank":
                answer_content = f"{quiz['answer']}"

            answer_slide.shapes.placeholders[1].text = answer_content

    prs.save(output_path)
    return output_path